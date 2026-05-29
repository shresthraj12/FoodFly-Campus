import os
import sys
import subprocess

# Auto-install python-pptx if missing
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("python-pptx not found. Installing...")
    subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette - Premium Light Theme
    bg_color = RGBColor(255, 255, 255)         # #FFFFFF (White Slide Canvas)
    slide_bg_cream = RGBColor(255, 247, 242)   # #FFF7F2 (Soft Cream sections)
    card_border_peach = RGBColor(255, 231, 219) # #FFE7DB (Peach border highlights)
    orange_accent = RGBColor(255, 107, 0)      # #FF6B00 (Primary Orange Accent)
    red_accent = RGBColor(255, 59, 48)         # #FF3B30 (Zomato/Blinkit Red Accent)
    
    # Typography Colors
    text_title = RGBColor(17, 24, 39)          # #111827 (Dark Charcoal Title Text)
    text_body = RGBColor(55, 65, 81)           # #374151 (Dark Gray Body Text)
    text_muted = RGBColor(107, 114, 128)       # #6B7280 (Slate Muted Text)
    white_pure = RGBColor(255, 255, 255)       # #FFFFFF (Pure White)

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        # Clean white canvas
        fill.fore_color.rgb = bg_color

    def add_slide_header(slide, slide_num_str, title_text):
        # Badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(0.4), Inches(0.8), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = slide_bg_cream
        badge.line.color.rgb = card_border_peach
        badge.line.width = Pt(1.5)
        tf_b = badge.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = 0
        p_b = tf_b.paragraphs[0]
        p_b.text = slide_num_str
        p_b.alignment = PP_ALIGN.CENTER
        p_b.font.name = 'Helvetica'
        p_b.font.size = Pt(13)
        p_b.font.bold = True
        p_b.font.color.rgb = orange_accent
        
        # Slide Title
        txBox = slide.shapes.add_textbox(Inches(1.75), Inches(0.35), Inches(10.833), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Helvetica'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = text_title

    def draw_card(slide, left, top, width, height, title, lines, accent=orange_accent, fill=white_pure, custom_title_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = accent
        shape.line.width = Pt(1.5)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Helvetica'
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = custom_title_color if custom_title_color else text_title
        p.space_after = Pt(8)
        
        for line in lines:
            p2 = tf.add_paragraph()
            p2.text = line
            p2.font.name = 'Helvetica'
            p2.font.size = Pt(11)
            p2.font.color.rgb = text_body
            p2.space_before = Pt(3)

    def try_add_image(slide, path, left, top, width, height, placeholder_title="Image"):
        if os.path.exists(path):
            try:
                slide.shapes.add_picture(path, left, top, width, height)
                return True
            except Exception as e:
                print(f"Error adding image {path}: {e}")
        
        # Fallback card
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = slide_bg_cream
        shape.line.color.rgb = card_border_peach
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[{placeholder_title}]\n{path}"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10)
        p.font.color.rgb = text_muted
        return False

    assets_dir = os.path.join("presentation", "assets")

    # ==================== SLIDE 1: Cover Slide ====================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    # Left Title Splash
    title_box = slide1.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(7.0), Inches(4.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p_badge = tf1.paragraphs[0]
    p_badge.text = "🔥 DevOps Case Study"
    p_badge.font.name = 'Helvetica'
    p_badge.font.size = Pt(13)
    p_badge.font.bold = True
    p_badge.font.color.rgb = orange_accent
    p_badge.space_after = Pt(16)
    
    p_title = tf1.add_paragraph()
    p_title.text = "FoodFly Campus"
    p_title.font.name = 'Helvetica'
    p_title.font.size = Pt(56)
    p_title.font.bold = True
    p_title.font.color.rgb = text_title
    
    p_sub = tf1.add_paragraph()
    p_sub.text = "Smart Campus Food Delivery Platform"
    p_sub.font.name = 'Helvetica'
    p_sub.font.size = Pt(24)
    p_sub.font.color.rgb = orange_accent
    p_sub.space_after = Pt(16)
    
    p_divider = tf1.add_paragraph()
    p_divider.text = "▬▬▬▬▬▬▬▬▬▬▬▬"
    p_divider.font.name = 'Helvetica'
    p_divider.font.size = Pt(14)
    p_divider.font.color.rgb = red_accent
    p_divider.space_after = Pt(16)
    
    p_meta = tf1.add_paragraph()
    p_meta.text = "A high-availability, multi-container student delivery platform engineered for LPU hostel lobbies. Integrated with React, Node, Mongo, Docker Compose, Jenkins CI/CD, and real-time observability telemetry."
    p_meta.font.name = 'Helvetica'
    p_meta.font.size = Pt(12)
    p_meta.font.color.rgb = text_muted

    # Right Cover Image
    try_add_image(slide1, os.path.join(assets_dir, "cover_illustration_light.png"), 
                  Inches(8.2), Inches(1.5), Inches(4.5), Inches(4.5), "Rider Illustration")

    # ==================== SLIDE 2: Problem Statement ====================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_slide_header(slide2, "02", "The Problem Statement")

    # Left: Problems list
    draw_card(slide2, Inches(0.75), Inches(1.8), Inches(6.0), Inches(1.5),
              "📍 Hostel Gate Restrictions",
              ["External delivery riders are restricted at LPU security checkpoints.", 
               "Students are forced to walk up to 1.5 km just to receive food packages."],
              accent=red_accent, fill=slide_bg_cream)
              
    draw_card(slide2, Inches(0.75), Inches(3.5), Inches(6.0), Inches(1.5),
              "⏱️ Long Delivery & Cold Food",
              ["Riders lack campus path coordinates, leading to routing confusion.", 
               "Prolonged transit times deliver cold, unappetizing meals to hostels."],
              accent=red_accent, fill=slide_bg_cream)
              
    draw_card(slide2, Inches(0.75), Inches(5.2), Inches(6.0), Inches(1.5),
              "🛒 Fragmented Campus Outlets",
              ["No digital platform aggregates individual tuck shops and canteens.", 
               "Students cannot checkout multiple vendors in a single transaction."],
              accent=red_accent, fill=slide_bg_cream)

    # Right: Problem Vector
    try_add_image(slide2, os.path.join(assets_dir, "problem_illustration_light.png"), 
                  Inches(7.3), Inches(1.8), Inches(5.3), Inches(4.9), "Problem Illustration")

    # ==================== SLIDE 3: Solution Overview ====================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_slide_header(slide3, "03", "The Solution Overview")

    # Left Image
    try_add_image(slide3, os.path.join(assets_dir, "solution_illustration_light.png"), 
                  Inches(0.75), Inches(1.8), Inches(5.3), Inches(4.9), "Solution Illustration")

    # Right: Solution Details
    draw_card(slide3, Inches(6.6), Inches(1.8), Inches(6.0), Inches(4.9),
              "Hyper-Local Campus Logistics Ecosystem",
              [
                  "FoodFly Campus resolves campus restrictions by creating a dedicated, closed-loop delivery network designed specifically for the university premises.",
                  "",
                  "⚡ Direct Hostel Block Deliveries:",
                  "  Riders deliver directly to designated tables at BH-1, BH-2, GH-1, etc.",
                  "",
                  "⚡ Authorized Campus Riders:",
                  "  Optimized internal routing using student riders who have gate access.",
                  "",
                  "⚡ DevOps Centered Resilience:",
                  "  Containerized stack built to auto-scale during peak dining hours."
              ],
              accent=orange_accent, fill=slide_bg_cream)

    # ==================== SLIDE 4: Key Features ====================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_slide_header(slide4, "04", "Key Platform Features")

    # Grid 3x2
    cw = Inches(3.6)
    ch = Inches(2.1)
    
    draw_card(slide4, Inches(0.75), Inches(1.8), cw, ch, "📍 Hostel Drop Zones", ["Direct delivery destinations tailored block-wise for BH-1 to BH-5, GH-1 to GH-3."])
    draw_card(slide4, Inches(4.85), Inches(1.8), cw, ch, "🍔 Multi-Vendor Cart", ["Checkout multiple tuck shops and canteens in a single payment transaction."])
    draw_card(slide4, Inches(8.95), Inches(1.8), cw, ch, "🚴 Rider Live Tracker", ["Step-by-step progress tracking from kitchen prep to packer hostel arrival."])
    
    draw_card(slide4, Inches(0.75), Inches(4.3), cw, ch, "🔑 Role Verification", ["Secure JWT checks with separate portals for Student, Seller, and DevOps Admin."])
    draw_card(slide4, Inches(4.85), Inches(4.3), cw, ch, "📈 System Observability", ["Integrated scrapers feeding CPU loads and Express API latencies to Prometheus."])
    draw_card(slide4, Inches(8.95), Inches(4.3), cw, ch, "🔄 Automated Rollouts", ["Webhook-triggered pipeline rebuilding docker images for zero-downtime deployments."])

    # ==================== SLIDE 5: Tech Stack ====================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_slide_header(slide5, "05", "Technical Stack Breakdown")

    # 4 Columns
    col_w = Inches(2.7)
    col_h = Inches(4.8)
    
    draw_card(slide5, Inches(0.75), Inches(1.8), col_w, col_h, "Frontend Tier",
              [
                  "• React.js SPA",
                  "• Tailwind CSS (v4)",
                  "• Vite Toolchain",
                  "• React Router DOM",
                  "• Context State API",
                  "• Axios HTTP Client",
                  "",
                  "Responsive styling mapping viewport configurations across mobile and desktop devices."
              ], accent=RGBColor(2, 132, 199)) # React Blue

    draw_card(slide5, Inches(3.75), Inches(1.8), col_w, col_h, "Backend API Tier",
              [
                  "• Node.js Runtime",
                  "• Express.js Framework",
                  "• JWT Token Auth",
                  "• RESTful Endpoints",
                  "• prom-client Telemetry",
                  "• CORS Middleware",
                  "",
                  "Asynchronous event-loop processes handling role-based routing checks."
              ], accent=RGBColor(67, 133, 61)) # Node Green

    draw_card(slide5, Inches(6.75), Inches(1.8), col_w, col_h, "Database Tier",
              [
                  "• MongoDB Community",
                  "• Mongoose ODM",
                  "• Named Docker Volumes",
                  "• BSON Data Formats",
                  "• Strict Schema Model",
                  "• Mapped DB Indexes",
                  "",
                  "NoSQL document store preserving delivery states across docker containers."
              ], accent=RGBColor(19, 170, 82)) # Mongo Green

    draw_card(slide5, Inches(9.75), Inches(1.8), col_w, col_h, "DevOps & Obs",
              [
                  "• Docker Containers",
                  "• Docker Compose",
                  "• Jenkins Pipelines",
                  "• Nginx Reverse Proxy",
                  "• Prometheus Collector",
                  "• Grafana Dashboard",
                  "",
                  "Git hook automation managing build cycles and serving telemetry graphs."
              ], accent=orange_accent)

    # ==================== SLIDE 6: System Architecture ====================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_slide_header(slide6, "06", "System Architecture & Flow")

    # Draw nodes with white/light backgrounds
    def draw_node_box(slide, x, y, w, h, text, color=white_pure, line_color=orange_accent):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Helvetica'
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = text_title

    # Draw nodes
    draw_node_box(slide6, Inches(0.75), Inches(2.2), Inches(1.8), Inches(0.8), "Client Browser\n(React UI)", line_color=RGBColor(2, 132, 199))
    draw_node_box(slide6, Inches(3.2), Inches(2.2), Inches(1.8), Inches(0.8), "Nginx Proxy\n(Port 8080)", line_color=RGBColor(16, 185, 129))
    draw_node_box(slide6, Inches(5.65), Inches(1.7), Inches(1.9), Inches(0.8), "React Container\n(Port 3000)", line_color=RGBColor(97, 218, 251))
    draw_node_box(slide6, Inches(5.65), Inches(2.7), Inches(1.9), Inches(0.8), "Express Container\n(Port 5000)", line_color=RGBColor(56, 189, 248))
    draw_node_box(slide6, Inches(8.2), Inches(2.7), Inches(1.8), Inches(0.8), "MongoDB Stack\n(Port 27017)", line_color=RGBColor(19, 170, 82))

    draw_node_box(slide6, Inches(0.75), Inches(4.6), Inches(1.8), Inches(0.8), "Jenkins Server\n(Port 8081)", color=slide_bg_cream, line_color=orange_accent)
    draw_node_box(slide6, Inches(3.2), Inches(4.6), Inches(1.8), Inches(0.8), "Docker Compose\n(Microservices)", color=slide_bg_cream, line_color=orange_accent)
    draw_node_box(slide6, Inches(5.65), Inches(4.6), Inches(1.9), Inches(0.8), "Node / cAdvisor\nExporters", color=slide_bg_cream, line_color=orange_accent)
    draw_node_box(slide6, Inches(8.2), Inches(4.6), Inches(1.8), Inches(0.8), "Prometheus DB\n(Port 9090)", color=slide_bg_cream, line_color=orange_accent)
    draw_node_box(slide6, Inches(10.7), Inches(4.6), Inches(1.8), Inches(0.8), "Grafana Dashboard\n(Port 3001)", color=slide_bg_cream, line_color=orange_accent)

    # Label Text boxes
    desc_box = slide6.shapes.add_textbox(Inches(0.75), Inches(6.0), Inches(11.833), Inches(0.9))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "Architecture Flow: Clients route to unified port 8080. Nginx proxies requests to static assets (Port 3000) or Express API endpoints (Port 5000) which read MongoDB. Git triggers webhook commits to Jenkins which builds and rolls out microservices via Docker Compose. cAdvisor and Node Exporter collect metrics scraped by Prometheus and graphed in Grafana."
    p_desc.font.name = 'Helvetica'
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = text_muted

    # ==================== SLIDE 7: Docker & Containerization ====================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_slide_header(slide7, "07", "Docker & Containerization")

    # Left: Code panel (with warm light tones)
    draw_card(slide7, Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.9),
              "Orchestration Setup (docker-compose.yml)",
              [
                  "services:",
                  "  mongodb: { image: mongo:latest, volumes: [mongodb_data] }",
                  "  backend:",
                  "    build: ./backend",
                  "    depends_on: { mongodb: { condition: healthy } }",
                  "    environment: [MONGODB_URI, JWT_SECRET]",
                  "  frontend:",
                  "    build: ./frontend",
                  "    ports: ['3000:80']",
                  "  nginx:",
                  "    image: nginx:alpine",
                  "    ports: ['8080:80']",
                  "    volumes: [./nginx/nginx.conf:/etc/nginx/nginx.conf]"
              ], accent=orange_accent, fill=white_pure)

    # Right Image top
    try_add_image(slide7, os.path.join(assets_dir, "docker_containers_light.png"),
                  Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.4), "Docker Illustration")
                  
    # Right Text bottom
    draw_card(slide7, Inches(6.8), Inches(4.4), Inches(5.7), Inches(2.3),
              "Containerization Highlights",
              [
                  "• Bridge Isolation: Services share 'foodfly_network' for private hostname resolution.",
                  "• Stage builds: Frontend reduces asset bundle weight from 350MB to 22MB.",
                  "• Persistent Volumes: Named volumes store MongoDB files across container cycles."
              ], accent=orange_accent, fill=slide_bg_cream)

    # ==================== SLIDE 8: Jenkins CI/CD Pipeline ====================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_slide_header(slide8, "08", "Jenkins CI/CD Build Pipeline")

    # Left: Pipeline stages
    draw_card(slide8, Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.9),
              "Declarative Automations (Jenkinsfile)",
              [
                  "⚙️ Webhook Triggering: GitHub hook polling automatically triggers builds on code pushes.",
                  "",
                  "⚙️ Pipeline Stage Logs:",
                  "  - Stage 1: Checkout Git code (5s)",
                  "  - Stage 2: Install node package dependencies (12s)",
                  "  - Stage 3: Build static production assets via Vite (45s)",
                  "  - Stage 4: Scan backend node files for syntax integrity (8s)",
                  "  - Stage 5: Compile multi-container Docker images (30s)",
                  "  - Stage 6: Spin up containers on target server (15s)",
                  "",
                  "⚙️ Workspace Cleanup: Prunes unused Docker assets automatically post-build."
              ], accent=orange_accent, fill=white_pure)

    # Right top: Jenkins Flow graphic
    try_add_image(slide8, os.path.join(assets_dir, "jenkins_illustration_light.png"),
                  Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.35), "Jenkins Illustration")

    # Right bottom: Jenkins Screenshot
    try_add_image(slide8, os.path.join(assets_dir, "jenkins_pipeline.png"),
                  Inches(6.8), Inches(4.35), Inches(5.7), Inches(2.35), "Jenkins Pipeline Screenshot")

    # ==================== SLIDE 9: Monitoring with Grafana & Prometheus ====================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_slide_header(slide9, "09", "Monitoring with Grafana & Prometheus")

    # Left top: Monitoring graphic (Reused)
    try_add_image(slide9, os.path.join(assets_dir, "grafana_observability.png"),
                  Inches(0.75), Inches(1.8), Inches(5.5), Inches(2.35), "Grafana Illustration")

    # Left bottom: Grafana Screenshot (Reused)
    try_add_image(slide9, os.path.join(assets_dir, "grafana_metrics.png"),
                  Inches(0.75), Inches(4.35), Inches(5.5), Inches(2.35), "Grafana Dashboard Screenshot")

    # Right: Telemetry Pillars
    draw_card(slide9, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.9),
              "Three-Dimensional Telemetry Scraping",
              [
                  "Scrapes target endpoints every 15s to monitor status & performance:",
                  "",
                  "📊 Application Tier (prom-client):",
                  "  Tracks Express request latency histogram thresholds and response HTTP status error rates.",
                  "",
                  "📊 Container Tier (cAdvisor):",
                  "  Tracks CPU percentage profiles, memory consumption limits, and disk I/O rates for each microservice container.",
                  "",
                  "📊 System Tier (Node Exporter):",
                  "  Gathers hardware levels (CPU load, RAM usage) of the root host VM."
              ], accent=orange_accent, fill=slide_bg_cream)

    # ==================== SLIDE 10: UI Screenshots & Conclusion (LIGHT REDESIGNED) ====================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_slide_header(slide10, "10", "UI Showcases & Pitch Summary")

    # Left: Conclusion points in modern icon cards (top part, y from 1.3 to 4.3)
    draw_card(slide10, Inches(0.75), Inches(1.4), Inches(5.8), Inches(3.0),
              "Key Takeaways & Startup Resiliency",
              [
                  "🐳 Microservice Portability:",
                  "  Docker containerization resolves environmental runtime discrepancies.",
                  "🔄 Automated Pipelines:",
                  "  Jenkins gates build rollouts if tests or compile steps fail.",
                  "📊 Operational Telemetry:",
                  "  Immediate status indicators on Grafana panels minimize downtime."
              ], accent=orange_accent, fill=white_pure)

    # Right: 4 Screenshots in a 2x2 grid (top part, y from 1.3 to 4.3)
    screenshots = [
        ("home_page.png", Inches(6.9), Inches(1.4)),
        ("menu_page.png", Inches(9.8), Inches(1.4)),
        ("checkout_page.png", Inches(6.9), Inches(2.95)),
        ("admin_dashboard.png", Inches(9.8), Inches(2.95))
    ]
    
    sc_w = Inches(2.7)
    sc_h = Inches(1.45)
    
    for filename, x, y in screenshots:
        border = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, sc_w, sc_h)
        border.fill.solid()
        border.fill.fore_color.rgb = white_pure
        border.line.color.rgb = card_border_peach
        border.line.width = Pt(1.5)
        
        img_p = os.path.join(assets_dir, filename)
        if os.path.exists(img_p):
            try:
                slide10.shapes.add_picture(img_p, x + Inches(0.04), y + Inches(0.04), sc_w - Inches(0.08), sc_h - Inches(0.08))
            except Exception as e:
                print(f"Error drawing screenshot {filename} on slide 10: {e}")
        else:
            tf_b = border.text_frame
            tf_b.word_wrap = True
            p_b = tf_b.paragraphs[0]
            p_b.text = f"[{filename}]"
            p_b.alignment = PP_ALIGN.CENTER
            p_b.font.size = Pt(10)
            p_b.font.color.rgb = text_muted

    # Bottom: Attractive "Thank You" block + footer icons + LPU branding (y from 4.6 to 6.9)
    thankyou_bg = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(4.7), Inches(11.833), Inches(2.2))
    thankyou_bg.fill.solid()
    thankyou_bg.fill.fore_color.rgb = slide_bg_cream
    thankyou_bg.line.color.rgb = card_border_peach
    thankyou_bg.line.width = Pt(1.5)

    # Thank you text in a transparent text box
    ty_box = slide10.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(6.0), Inches(2.0))
    tf_ty = ty_box.text_frame
    tf_ty.word_wrap = True
    
    p_ty = tf_ty.paragraphs[0]
    p_ty.text = "Thank You!"
    p_ty.font.name = 'Helvetica'
    p_ty.font.size = Pt(36)
    p_ty.font.bold = True
    p_ty.font.color.rgb = orange_accent
    
    p_ty_sub = tf_ty.add_paragraph()
    p_ty_sub.text = "FoodFly Campus — Revolutionizing Dining Logistics at LPU"
    p_ty_sub.font.name = 'Helvetica'
    p_ty_sub.font.size = Pt(13)
    p_ty_sub.font.bold = True
    p_ty_sub.font.color.rgb = text_body
    p_ty_sub.space_after = Pt(12)

    # Startup Footer icons (Drawn as small text inside pill shapes)
    badges_info = [
        ("🚴 Fast Delivery", Inches(1.0)),
        ("🍔 Fresh Food", Inches(2.4)),
        ("🏫 Campus First", Inches(3.7)),
        ("🎓 Student Focus", Inches(5.1))
    ]
    for text, badge_x in badges_info:
        pill = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, badge_x, Inches(6.1), Inches(1.25), Inches(0.35))
        pill.fill.solid()
        pill.fill.fore_color.rgb = white_pure
        pill.line.color.rgb = card_border_peach
        pill.line.width = Pt(1)
        tf_p = pill.text_frame
        tf_p.word_wrap = True
        tf_p.margin_left = tf_p.margin_right = tf_p.margin_top = tf_p.margin_bottom = 0
        p_p = tf_p.paragraphs[0]
        p_p.text = text
        p_p.alignment = PP_ALIGN.CENTER
        p_p.font.name = 'Helvetica'
        p_p.font.size = Pt(9)
        p_p.font.bold = True
        p_p.font.color.rgb = text_title

    # LPU Text/Branding card on bottom right of Slide 10
    lpu_brand_card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(4.9), Inches(4.5), Inches(1.8))
    lpu_brand_card.fill.solid()
    lpu_brand_card.fill.fore_color.rgb = white_pure
    lpu_brand_card.line.color.rgb = orange_accent
    lpu_brand_card.line.width = Pt(1.5)
    tf_lpu = lpu_brand_card.text_frame
    tf_lpu.word_wrap = True
    tf_lpu.margin_left = Inches(0.2)
    tf_lpu.margin_right = Inches(0.2)
    tf_lpu.margin_top = Inches(0.2)
    p_lpu_title = tf_lpu.paragraphs[0]
    p_lpu_title.text = "🎓 Designed for LPU Campus"
    p_lpu_title.font.name = 'Helvetica'
    p_lpu_title.font.size = Pt(14)
    p_lpu_title.font.bold = True
    p_lpu_title.font.color.rgb = red_accent
    p_lpu_title.space_after = Pt(6)
    
    p_lpu_desc = tf_lpu.add_paragraph()
    p_lpu_desc.text = "Hostel delivery zone mappings customized for LPU blocks (BH-1 to BH-6, GH-1 to GH-3) supporting local vendors and security permissions."
    p_lpu_desc.font.name = 'Helvetica'
    p_lpu_desc.font.size = Pt(10)
    p_lpu_desc.font.color.rgb = text_body

    # Save presentation
    output_path = os.path.join("presentation", "FoodFly_Campus_Light_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation successfully compiled to: {output_path}")

if __name__ == '__main__':
    create_presentation()
